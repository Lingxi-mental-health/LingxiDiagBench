# -*- coding: utf-8 -*-
"""
文档加载模块
支持多种文件格式：PDF, Word, PPT, Excel, TXT, JSON
"""

import json
from pathlib import Path
from typing import List, Dict, Any
import logging

# 文档处理库
try:
    from pypdf import PdfReader
except ImportError:
    PdfReader = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import openpyxl
    import pandas as pd
except ImportError:
    openpyxl = None
    pd = None

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class DocumentLoader:
    """文档加载器，支持多种文件格式"""
    
    def __init__(self):
        """初始化文档加载器"""
        self.supported_formats = {
            '.pdf': self._load_pdf,
            '.docx': self._load_docx,
            '.doc': self._load_docx,
            '.pptx': self._load_pptx,
            '.ppt': self._load_pptx,
            '.xlsx': self._load_excel,
            '.xls': self._load_excel,
            '.txt': self._load_txt,
            '.json': self._load_json,
        }
    
    def load_document(self, file_path: str) -> Dict[str, Any]:
        """
        加载文档
        
        Args:
            file_path: 文件路径
            
        Returns:
            包含文档内容和元数据的字典
        """
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise FileNotFoundError(f"文件不存在: {file_path}")
        
        suffix = file_path.suffix.lower()
        
        if suffix not in self.supported_formats:
            raise ValueError(
                f"不支持的文件格式: {suffix}. "
                f"支持的格式: {', '.join(self.supported_formats.keys())}"
            )
        
        logger.info(f"正在加载文档: {file_path.name}")
        
        # 调用对应的加载函数
        loader_func = self.supported_formats[suffix]
        content = loader_func(file_path)
        
        # 构建文档结构
        document = {
            'file_name': file_path.name,
            'file_path': str(file_path.absolute()),
            'file_type': suffix,
            'content': content,
            'metadata': {
                'size': file_path.stat().st_size,
                'modified_time': file_path.stat().st_mtime,
            }
        }
        
        logger.info(f"✓ 文档加载完成: {file_path.name}")
        return document
    
    def load_directory(self, dir_path: str, recursive: bool = True) -> List[Dict[str, Any]]:
        """
        加载目录下的所有文档
        
        Args:
            dir_path: 目录路径
            recursive: 是否递归加载子目录
            
        Returns:
            文档列表
        """
        dir_path = Path(dir_path)
        
        if not dir_path.exists():
            raise FileNotFoundError(f"目录不存在: {dir_path}")
        
        if not dir_path.is_dir():
            raise ValueError(f"路径不是目录: {dir_path}")
        
        documents = []
        
        # 获取文件列表
        if recursive:
            files = list(dir_path.rglob('*'))
        else:
            files = list(dir_path.glob('*'))
        
        # 过滤文件
        files = [f for f in files if f.is_file() and f.suffix.lower() in self.supported_formats]
        
        logger.info(f"找到 {len(files)} 个文档文件")
        
        # 加载每个文件
        for file_path in files:
            try:
                doc = self.load_document(str(file_path))
                documents.append(doc)
            except Exception as e:
                logger.error(f"加载文档失败 {file_path.name}: {str(e)}")
        
        logger.info(f"✓ 成功加载 {len(documents)} 个文档")
        return documents
    
    def _load_pdf(self, file_path: Path) -> str:
        """加载PDF文件"""
        if PdfReader is None:
            raise ImportError("需要安装pypdf库: pip install pypdf")
        
        reader = PdfReader(str(file_path))
        text_parts = []
        
        for page in reader.pages:
            text = page.extract_text()
            if text:
                text_parts.append(text)
        
        return '\n\n'.join(text_parts)
    
    def _load_docx(self, file_path: Path) -> str:
        """加载Word文件"""
        if Document is None:
            raise ImportError("需要安装python-docx库: pip install python-docx")
        
        doc = Document(str(file_path))
        text_parts = []
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)
        
        # 提取表格内容
        for table in doc.tables:
            for row in table.rows:
                row_text = ' | '.join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    text_parts.append(row_text)
        
        return '\n\n'.join(text_parts)
    
    def _load_pptx(self, file_path: Path) -> str:
        """加载PPT文件"""
        if Presentation is None:
            raise ImportError("需要安装python-pptx库: pip install python-pptx")
        
        prs = Presentation(str(file_path))
        text_parts = []
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)
        
        return '\n\n'.join(text_parts)
    
    def _load_excel(self, file_path: Path) -> str:
        """加载Excel文件"""
        if pd is None:
            raise ImportError("需要安装pandas和openpyxl库: pip install pandas openpyxl")
        
        # 读取所有sheet
        excel_file = pd.ExcelFile(file_path)
        text_parts = []
        
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(excel_file, sheet_name=sheet_name)
            
            # 添加sheet名称
            text_parts.append(f"=== {sheet_name} ===")
            
            # 转换为文本
            text_parts.append(df.to_string(index=False))
        
        return '\n\n'.join(text_parts)
    
    def _load_txt(self, file_path: Path) -> str:
        """加载TXT文件"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    
    def _load_json(self, file_path: Path) -> str:
        """加载JSON文件"""
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        
        # 将JSON转换为可读的文本格式
        return json.dumps(data, ensure_ascii=False, indent=2)

